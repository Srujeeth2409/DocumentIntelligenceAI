"""
conversion.py - File format conversion utilities (MacOS Compatible)
Handles conversions between PDF, Word, Excel, PowerPoint, and Images
Includes PDF merging and password protection
"""

import os
import sys
import tempfile
import subprocess
from PIL import Image
from pdf2image import convert_from_path
import pytesseract
from docx import Document
from docx.shared import Inches
import zipfile

# Handle different PyPDF2/pypdf versions
try:
    from pypdf import PdfMerger, PdfReader, PdfWriter

    print("[INFO] Using pypdf library")
except ImportError:
    try:
        from PyPDF2 import PdfMerger, PdfReader, PdfWriter

        print("[INFO] Using PyPDF2 3.0+")
    except ImportError:
        try:
            from PyPDF2 import PdfFileMerger as PdfMerger
            from PyPDF2 import PdfFileReader as PdfReader
            from PyPDF2 import PdfFileWriter as PdfWriter

            print("[INFO] Using PyPDF2 2.0")
        except ImportError:
            from PyPDF2 import PdfFileMerger as PdfMerger
            from PyPDF2 import PdfFileReader as PdfReader
            from PyPDF2 import PdfFileWriter as PdfWriter

            print("[INFO] Using legacy PyPDF2")


# =====================================================================
# IMAGE TO PDF
# =====================================================================

def jpg_to_pdf(input_path, output_path):
    """
    Convert JPG/PNG image to PDF.

    Args:
        input_path: Path to input image file
        output_path: Path to save PDF
    """
    try:
        image = Image.open(input_path)

        # Convert RGBA to RGB if necessary
        if image.mode == 'RGBA':
            image = image.convert('RGB')

        # Save as PDF
        image.save(output_path, 'PDF', resolution=100.0, quality=95)

        print(f"[SUCCESS] Image converted to PDF: {output_path}")
        return output_path

    except Exception as e:
        print(f"[ERROR] Image to PDF conversion failed: {e}")
        raise


# =====================================================================
# WORD TO PDF (MacOS - using LibreOffice)
# =====================================================================

def word_to_pdf(input_path, output_path):
    """
    Convert Word document to PDF using LibreOffice (MacOS).

    Args:
        input_path: Path to input Word file (.doc or .docx)
        output_path: Path to save PDF
    """
    try:
        # Ensure absolute paths
        input_path = os.path.abspath(input_path)
        output_dir = os.path.dirname(os.path.abspath(output_path))

        # Try using LibreOffice
        libreoffice_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/usr/local/bin/soffice',
            'soffice'
        ]

        soffice_path = None
        for path in libreoffice_paths:
            if os.path.exists(path) or path == 'soffice':
                soffice_path = path
                break

        if soffice_path:
            # Convert using LibreOffice
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_dir,
                input_path
            ]

            subprocess.run(cmd, check=True, capture_output=True)

            # LibreOffice saves with original filename, rename if needed
            base_name = os.path.splitext(os.path.basename(input_path))[0]
            generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")

            if generated_pdf != output_path and os.path.exists(generated_pdf):
                os.rename(generated_pdf, output_path)

            print(f"[SUCCESS] Word converted to PDF: {output_path}")
            return output_path
        else:
            # Fallback: Use python-docx to extract text and create PDF
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet

            doc_word = Document(input_path)

            # Create PDF
            pdf = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            for para in doc_word.paragraphs:
                if para.text.strip():
                    p = Paragraph(para.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))

            pdf.build(story)

            print(f"[SUCCESS] Word converted to PDF (text only): {output_path}")
            return output_path

    except Exception as e:
        print(f"[ERROR] Word to PDF conversion failed: {e}")
        raise


# =====================================================================
# POWERPOINT TO PDF (MacOS - using LibreOffice)
# =====================================================================

def ppt_to_pdf(input_path, output_path):
    """
    Convert PowerPoint presentation to PDF using LibreOffice (MacOS).

    Args:
        input_path: Path to input PowerPoint file (.ppt or .pptx)
        output_path: Path to save PDF
    """
    try:
        # Ensure absolute paths
        input_path = os.path.abspath(input_path)
        output_dir = os.path.dirname(os.path.abspath(output_path))

        # Try using LibreOffice
        libreoffice_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/usr/local/bin/soffice',
            'soffice'
        ]

        soffice_path = None
        for path in libreoffice_paths:
            if os.path.exists(path) or path == 'soffice':
                soffice_path = path
                break

        if soffice_path:
            # Convert using LibreOffice
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_dir,
                input_path
            ]

            subprocess.run(cmd, check=True, capture_output=True)

            # LibreOffice saves with original filename, rename if needed
            base_name = os.path.splitext(os.path.basename(input_path))[0]
            generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")

            if generated_pdf != output_path and os.path.exists(generated_pdf):
                os.rename(generated_pdf, output_path)

            print(f"[SUCCESS] PowerPoint converted to PDF: {output_path}")
            return output_path
        else:
            raise Exception("LibreOffice not found. Please install LibreOffice for PowerPoint to PDF conversion.")

    except Exception as e:
        print(f"[ERROR] PowerPoint to PDF conversion failed: {e}")
        raise


# =====================================================================
# EXCEL TO PDF (MacOS - using LibreOffice)
# =====================================================================

def excel_to_pdf(input_path, output_path):
    """
    Convert Excel spreadsheet to PDF using LibreOffice (MacOS).

    Args:
        input_path: Path to input Excel file (.xls or .xlsx)
        output_path: Path to save PDF
    """
    try:
        # Ensure absolute paths
        input_path = os.path.abspath(input_path)
        output_dir = os.path.dirname(os.path.abspath(output_path))

        # Try using LibreOffice
        libreoffice_paths = [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/usr/local/bin/soffice',
            'soffice'
        ]

        soffice_path = None
        for path in libreoffice_paths:
            if os.path.exists(path) or path == 'soffice':
                soffice_path = path
                break

        if soffice_path:
            # Convert using LibreOffice
            cmd = [
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_dir,
                input_path
            ]

            subprocess.run(cmd, check=True, capture_output=True)

            # LibreOffice saves with original filename, rename if needed
            base_name = os.path.splitext(os.path.basename(input_path))[0]
            generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")

            if generated_pdf != output_path and os.path.exists(generated_pdf):
                os.rename(generated_pdf, output_path)

            print(f"[SUCCESS] Excel converted to PDF: {output_path}")
            return output_path
        else:
            raise Exception("LibreOffice not found. Please install LibreOffice for Excel to PDF conversion.")

    except Exception as e:
        print(f"[ERROR] Excel to PDF conversion failed: {e}")
        raise


# =====================================================================
# PDF TO IMAGE
# =====================================================================

def pdf_to_jpg(input_path, output_folder):
    """
    Convert PDF pages to JPG images.

    Args:
        input_path: Path to input PDF file
        output_folder: Folder to save JPG images

    Returns:
        List of paths to generated JPG files
    """
    try:
        # Convert PDF to images
        images = convert_from_path(input_path, dpi=300)

        jpg_files = []
        for i, image in enumerate(images, start=1):
            output_path = os.path.join(output_folder, f'page_{i}.jpg')
            image.save(output_path, 'JPEG', quality=95)
            jpg_files.append(output_path)

        print(f"[SUCCESS] PDF converted to {len(jpg_files)} JPG images")
        return jpg_files

    except Exception as e:
        print(f"[ERROR] PDF to JPG conversion failed: {e}")
        raise


# =====================================================================
# PDF TO WORD
# =====================================================================

def pdf_to_word(input_path, output_path):
    """
    Convert PDF to Word document by extracting text and images.

    Args:
        input_path: Path to input PDF file
        output_path: Path to save Word document
    """
    try:
        # Convert PDF pages to images
        images = convert_from_path(input_path, dpi=200)

        # Create Word document
        doc = Document()

        for i, image in enumerate(images, start=1):
            # Extract text using OCR
            text = pytesseract.image_to_string(image)

            # Add text to document
            if text.strip():
                doc.add_paragraph(text)

            # Add page image
            temp_image_path = os.path.join(tempfile.gettempdir(), f'page_{i}.png')
            image.save(temp_image_path, 'PNG')

            try:
                doc.add_picture(temp_image_path, width=Inches(6))
            except:
                pass  # Skip if image can't be added

            # Add page break (except for last page)
            if i < len(images):
                doc.add_page_break()

        # Save document
        doc.save(output_path)

        print(f"[SUCCESS] PDF converted to Word: {output_path}")
        return output_path

    except Exception as e:
        print(f"[ERROR] PDF to Word conversion failed: {e}")
        raise


# =====================================================================
# PDF TO POWERPOINT
# =====================================================================

def pdf_to_ppt(input_path, output_path):
    """
    Convert PDF to PowerPoint by converting pages to images.

    Args:
        input_path: Path to input PDF file
        output_path: Path to save PowerPoint file
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches

        # Convert PDF pages to images
        images = convert_from_path(input_path, dpi=150)

        # Create presentation
        prs = Presentation()

        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for i, image in enumerate(images, start=1):
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Save image temporarily
            temp_image_path = os.path.join(tempfile.gettempdir(), f'page_{i}.png')
            image.save(temp_image_path, 'PNG')

            # Add image to slide (centered)
            left = Inches(0.5)
            top = Inches(0.5)
            height = Inches(6.5)

            slide.shapes.add_picture(temp_image_path, left, top, height=height)

        # Save presentation
        prs.save(output_path)

        print(f"[SUCCESS] PDF converted to PowerPoint: {output_path}")
        return output_path

    except Exception as e:
        print(f"[ERROR] PDF to PowerPoint conversion failed: {e}")
        raise


# =====================================================================
# PDF TO EXCEL
# =====================================================================

def pdf_to_excel(input_path, output_path):
    """
    Convert PDF to Excel by extracting text and attempting to parse tables.

    Args:
        input_path: Path to input PDF file
        output_path: Path to save Excel file
    """
    try:
        import pandas as pd
        from openpyxl import Workbook

        try:
            import tabula

            # Extract tables from PDF
            tables = tabula.read_pdf(input_path, pages='all', multiple_tables=True)

            if tables and len(tables) > 0:
                # Save all tables to Excel (each table on a separate sheet)
                with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                    for i, table in enumerate(tables, start=1):
                        sheet_name = f'Table_{i}'
                        table.to_excel(writer, sheet_name=sheet_name, index=False)

                print(f"[SUCCESS] PDF converted to Excel: {output_path}")
                return output_path
        except Exception as tabula_error:
            print(f"[WARNING] Tabula extraction failed: {tabula_error}")

        # Fallback: Extract text and create simple Excel
        images = convert_from_path(input_path, dpi=150)
        extracted_text = []

        for image in images:
            text = pytesseract.image_to_string(image)
            extracted_text.append(text)

        # Create Excel workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted Text"

        # Add headers
        ws.cell(row=1, column=1, value="Page")
        ws.cell(row=1, column=2, value="Content")

        # Write text to cells (one page per row)
        for i, text in enumerate(extracted_text, start=1):
            ws.cell(row=i + 1, column=1, value=f"Page {i}")
            ws.cell(row=i + 1, column=2, value=text)

        wb.save(output_path)

        print(f"[SUCCESS] PDF converted to Excel (text extraction): {output_path}")
        return output_path

    except Exception as e:
        print(f"[ERROR] PDF to Excel conversion failed: {e}")
        raise


# =====================================================================
# MERGE PDFs
# =====================================================================

def merge_pdfs(pdf_paths, output_path):
    """
    Merge multiple PDF files into one.

    Args:
        pdf_paths: List of paths to PDF files to merge
        output_path: Path where merged PDF will be saved
    """
    try:
        merger = PdfMerger()

        for pdf_path in pdf_paths:
            print(f"[INFO] Adding {pdf_path} to merge...")
            merger.append(pdf_path)

        # Write merged PDF
        merger.write(output_path)
        merger.close()

        print(f"[SUCCESS] Merged {len(pdf_paths)} PDFs into {output_path}")
        return output_path

    except Exception as e:
        print(f"[ERROR] PDF merge failed: {e}")
        raise


# =====================================================================
# PASSWORD PROTECT PDF
# =====================================================================

def protect_pdf(input_path, output_path, password):
    """
    Add password protection to a PDF file.

    Args:
        input_path: Path to input PDF
        output_path: Path to save protected PDF
        password: Password to protect the PDF
    """
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()

        # Check if using old or new API
        try:
            # Try new API first
            for page in reader.pages:
                writer.add_page(page)

            # Try encrypting with new API parameters
            try:
                writer.encrypt(user_password=password, owner_password=password, algorithm="AES-256")
            except TypeError:
                writer.encrypt(user_password=password, owner_password=password)

        except AttributeError:
            # Fall back to old API
            for page_num in range(reader.getNumPages()):
                writer.addPage(reader.getPage(page_num))

            writer.encrypt(user_pwd=password, owner_pwd=password)

        # Write protected PDF
        with open(output_path, "wb") as output_file:
            writer.write(output_file)

        print(f"[SUCCESS] PDF protected with password: {output_path}")
        return output_path

    except Exception as e:
        print(f"[ERROR] PDF protection failed: {e}")
        raise


# =====================================================================
# HELPER FUNCTIONS
# =====================================================================

def get_file_extension(filename):
    """Get file extension without dot."""
    return filename.split('.')[-1].lower()


def is_valid_conversion(source_ext, target_format):
    """
    Check if conversion is valid.
    """
    valid_conversions = {
        'pdf': ['jpg', 'word', 'ppt', 'excel'],
        'jpg': ['pdf'],
        'jpeg': ['pdf'],
        'png': ['pdf'],
        'doc': ['pdf'],
        'docx': ['pdf'],
        'ppt': ['pdf'],
        'pptx': ['pdf'],
        'xls': ['pdf'],
        'xlsx': ['pdf']
    }

    return target_format in valid_conversions.get(source_ext, [])


if __name__ == "__main__":
    print("Conversion utilities loaded successfully (MacOS)")